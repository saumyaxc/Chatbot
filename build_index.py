# build_index.py
# - Read every supported file in the docs/ folder (pptx, docx, pdf-text, txt, md)
# - Extract human-readable text (slide-by-slide, page-by-page for pptx/pdf)
# - Split long text into smaller chunks (better for search quality)
# - Create embeddings (numeric fingerprints) for each chunk
# - Save:
#     index/embeddings.npy  -> matrix of embeddings
#     index/meta.json       -> list of {source, text} aligned 1:1 with embeddings

import os, json, math
import numpy as np
from dotenv import load_dotenv
from openai import OpenAI

# Readers for different file types
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document
from pypdf import PdfReader

load_dotenv()   # loads API keys from .env

DOCS_DIR = "docs"
INDEX_DIR = "index"
EMBED_MODEL = "text-embedding-3-small"

client = OpenAI(api_key=os.environ["OPENAI_API_KEY"])

# Read plain text files.
def read_txt(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

# Read Markdown (treated as plain text).
def read_md(path):
    return read_txt(path)

# Read DOCX; join paragraph texts.
def read_docx(path):
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs)


# Read text from a PDF file.
# Returns list of (source_label, text) tuples, one per page.
# Example source_label: "file.pdf#page-3"
def read_pdf(path):
    reader = PdfReader(path)
    texts = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            t = page.extract_text() or ""
        except Exception:
            t = ""
        if t.strip():
            texts.append((f"{os.path.basename(path)}#page-{i}", t))
    return texts

# Read PPTX.
# - Captures slide titles, text boxes, tables, grouped-shape text, and slide notes.
# - Skips things without text (charts, connectors, images).
# Returns list of (source_label, text) per slide, e.g. 'Deck.pptx#slide-3'

def read_pptx(path):
    prs = Presentation(path)
    items = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = []

        # Title placeholder (if any)
        try:
            if getattr(slide.shapes, "title", None) and slide.shapes.title.has_text_frame:
                t = slide.shapes.title.text.strip()
                if t:
                    parts.append(t)
        except Exception:
            pass

        # All other shapes
        for shape in slide.shapes:
            # Skip the title we already captured
            if shape is getattr(slide.shapes, "title", None):
                continue

            try:
                # Text boxes / placeholders
                if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                    t = shape.text.strip()
                    if t:
                        parts.append(t)

                # Tables
                elif getattr(shape, "has_table", False) and shape.has_table:
                    tbl = shape.table
                    for row in tbl.rows:
                        cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
                        if cells:
                            parts.append(" | ".join(cells))

                # Grouped shapes: walk inside the group
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for shp in shape.shapes:
                        if getattr(shp, "has_text_frame", False) and shp.has_text_frame:
                            t = shp.text.strip()
                            if t:
                                parts.append(t)
                        elif getattr(shp, "has_table", False) and shp.has_table:
                            tbl = shp.table
                            for row in tbl.rows:
                                cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
                                if cells:
                                    parts.append(" | ".join(cells))
                # Everything else (charts, connectors, images…) → skip
            except Exception:
                # If any shape is weird, just skip it
                continue

        try:
            if getattr(slide, "notes_slide", None) and slide.notes_slide:
                n = slide.notes_slide.notes_text_frame.text
                if n and n.strip():
                    parts.append("Notes: " + n.strip())
        except Exception:
            pass

        text = "\n".join(p for p in parts if p)
        if text.strip():
            items.append((f"{os.path.basename(path)}#slide-{i}", text))

    return items  # list of (source, text)

# Split long text into overlapping chunks.
# - Smaller chunks search better and fit context windows.
# - Overlap helps keep sentences/ideas from being cut off abruptly.
def chunk_text(text, max_chars=1200, overlap=200):
    text = text.strip()
    if not text:
        return []
    chunks = []
    i = 0
    while i < len(text):
        chunk = text[i:i+max_chars]
        chunks.append(chunk)
        i += max_chars - overlap
    return chunks

# Create embeddings for a list of strings using the embedding API.
# Batch requests to reduce API calls
# Returns array of shape (len(texts), embedding_dim)
def embed_texts(texts):
    # texts: list[str]
    # batch to reduce API calls
    batch_size = 64 # adjust if needed
    out = []
    for i in range(0, len(texts), batch_size):
        batch = texts[i:i+batch_size]
        response = client.embeddings.create(model=EMBED_MODEL, input=batch)
        out.extend([d.embedding for d in response.data])
    return np.array(out, dtype=np.float32)

def main():
    os.makedirs(INDEX_DIR, exist_ok=True)
    records = []
    for root, _, files in os.walk(DOCS_DIR):
        for fname in files:
            path = os.path.join(root, fname)
            lower = fname.lower()
            try:
                if lower.endswith(".txt"):
                    text = read_txt(path)
                    for ch in chunk_text(text):
                        records.append({"source": os.path.basename(path), "text": ch})
                elif lower.endswith(".md"):
                    text = read_md(path)
                    for ch in chunk_text(text):
                        records.append({"source": os.path.basename(path), "text": ch})
                elif lower.endswith(".docx"):
                    text = read_docx(path)
                    for ch in chunk_text(text):
                        records.append({"source": os.path.basename(path), "text": ch})
                elif lower.endswith(".pdf"):
                    for src, t in read_pdf(path):
                        for ch in chunk_text(t):
                            records.append({"source": src, "text": ch})
                elif lower.endswith(".pptx"):
                    for src, t in read_pptx(path):
                        for ch in chunk_text(t):
                            records.append({"source": src, "text": ch})
                else:
                    print(f" (skip unsupported type) {fname}")
            except Exception as e:
                print(f" (error reading {fname}: {e})")
    if not records:
        print("No text found. Exiting.")
        return

    print(f"Embedding {len(records)} chunks...")
    texts = [r["text"] for r in records]
    embs = embed_texts(texts)

    norms = np.linalg.norm(embs, axis=1, keepdims=True) + 1e-12
    embs = embs / norms

    np.save(os.path.join(INDEX_DIR, "embeddings.npy"), embs)
    with open(os.path.join(INDEX_DIR, "meta.json"), "w", encoding="utf-8") as f:
        json.dump(records, f, ensure_ascii=False)
    

    print(f"Index built with {len(records)} chunks.")
    print(f"- {INDEX_DIR}/embeddings.npy")
    print(f"- {INDEX_DIR}/meta.json")

if __name__ == "__main__":
    main()